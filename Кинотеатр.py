"""программа подразумевает, что пользователь не делает ошибок при вводе"""
from pptx import Presentation
from docx import Document
import xlsxwriter

cnt_cinemas = 1
cnt_cinema_halls = 1
cnt_sessions = 1

cinemas = []
cinema_halls = []
sessions = []


class Cinema:
    def __init__(self, n, name, address):
        self.n = n
        self.name = name
        self.address = address

    def __str__(self):
        return f'Кинотеатр №{self.n} {self.name} по адресу:{self.address}'


class Hall:
    def __init__(self, n, cinema, chairs):
        self.n = n
        self.cinema = cinema
        self.chairs = chairs

    def __str__(self):
        return f'Зал №{self.n} кинотеатра {self.cinema}\n' + 'Конфигурация кресел:\n' + \
            ('  ' + ' '.join(str(i) for i in range(1, len(self.chairs) + 1))) + '\n' + \
            '\n'.join(str(i + 1) + ' ' + ' '.join(self.chairs[i]) for i in range(len(self.chairs)))


class Session:
    def __init__(self, n, name, cinema, hall, start_time, duration):
        self.n = n
        self.name = name
        self.cinema = cinema
        self.chairs = hall.chairs
        self.start_time = start_time
        self.duration = duration
        self.n_hall = hall.n
        self.text = f'Сеанс {self.n}, фильм {self.name} в зале № {self.n} кинотеатра {self.cinema}, начинающееся в ' + \
                    f'{self.start_time} ' + \
                    f'и длительностью {self.duration}\n'

    def __str__(self):
        return f'Сеанс {self.n}, фильм {self.name} в зале № {self.n} кинотеатра {self.cinema}, начинающееся в ' + \
                f'{self.start_time} ' + \
                f'и длительностью {self.duration}\n' + \
                ('  ' + ' '.join(str(i) for i in range(1, len(self.chairs) + 1))) + '\n' + \
                '\n'.join(str(i + 1) + ' ' + ' '.join(self.chairs[i]) for i in range(len(self.chairs)))


def main():
    global cnt_cinema_halls, cnt_sessions, cnt_cinemas

    while True:
        print()
        print('Выберите желаемое действие:')
        print('0 - выход')
        print('1 - добавить кинотеатр')
        print('2 - вывести список кинотеатров')
        print('3 - добавить зал')
        print('4 - вывести список залов')
        print('5 - добавить сеанс')
        print('6 - вывести список сеансов')
        print('7 - купить билет')
        print('8 - расписание фильмов в формате pptx')
        print('9 - рекламные буклеты в формате docx')
        print('10 - график загруженности кинотеатров в формате xlsx')

        n = int(input())

        if not n:
            break

        if n == 1:
            name = input('Название кинотеатра: ')
            address = input('Адрес кинотеатра: ')
            cinemas.append(Cinema(cnt_cinemas, name, address))
            cnt_cinemas += 1
            print('Кинотеатр успешно добавлен!')

        elif n == 2:
            for i in cinemas:
                print(i)

        elif n == 3:
            cinema = input('Название кинотеатра: ')
            r = int(input('Количество рядов в зале: '))
            c = int(input('Количество мест в 1м ряду: '))
            cinema_halls.append(Hall(cnt_cinema_halls, cinema, [['o' for _ in range(c)] for _ in range(r)]))
            cnt_cinema_halls += 1
            print('Зал успешно добавлен!')

        elif n == 4:
            for i in cinema_halls:
                print(i)

        elif n == 5:
            name = input('Название сеанса(фильма): ')
            cinema = input('Название кинотеатра: ')
            hall = cinema_halls[int(input('Номер зала: ')) - 1]
            start_time = input('Начало фильма(время): ')
            duration = input('Длительность фильма: ')
            sessions.append(Session(cnt_sessions, name, cinema, hall, start_time, duration))
            cnt_sessions += 1
            print('Сеанс успешно добавлен!')

        elif n == 6:
            for i in sessions:
                print(i)

        elif n == 7:
            for i in sessions:
                print(i)
            print('Введите номер сеанса, на который хотите приобрести билет')
            n = int(input()) - 1
            print(sessions[n])
            r = int(input('Номер ряда в зале: '))
            c = int(input('Номер места в ряду: '))
            sessions[n].chairs[r - 1][c - 1] = 'x'

            print('---------------------БИЛЕТ----------------------')
            print(f'Название кинотеатра: {sessions[n].cinema}')
            print(f'Название фильма: {sessions[n].cinema}')
            print(f'Номер зала: {sessions[n].n_hall}')
            print(f'Номер ряда: {r}')
            print(f'Номер места: {c}')
            print()
            print('Билет успешно куплен!')
            print('Приятного просмотра!')

        elif n == 8:
            # создаем новую презентацию
            prs = Presentation()
            # получаем схему расположения элементов для заголовочного слайда
            title_slide_layout = prs.slide_layouts[0]
            # создаем заголовочный слайд
            slide = prs.slides.add_slide(title_slide_layout)
            # создаем у слайда заголовок и текст
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "Расписание сеансов"
            subtitle.text = "На ближайшее время"
            for i in range(len(sessions)):
                # создаем новый слайд со схемой для добавления изображений
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = sessions[i].text
            # сохраняем презентацию
            prs.save('prs.pptx')

        elif n == 9:
            doc = Document()
            doc.add_heading('Рекламные билеты', 0)

            doc.add_page_break()
            for i in range(len(sessions) - 1):
                doc.add_heading('ТОЛЬКО У НАС!!!', 1)
                doc.add_paragraph(sessions[i].text, style='Intense Quote')
                doc.add_page_break()

            doc.add_heading('ТОЛЬКО У НАС!!!', 1)
            doc.add_paragraph(sessions[-1].text, style='Intense Quote')
            doc.save('test.docx')

        else:
            # Создаем новый Excel-файл и добавляем в него лист
            workbook = xlsxwriter.Workbook('data.xlsx')
            worksheet = workbook.add_worksheet()

            # Создаем данные в виде словаря
            data_dict = {i.name: [0] for i in cinemas}
            for i in sessions:
                data_dict[i.cinema][0] += 1

            # Записываем заголовки
            worksheet.write_row(0, 0, ['Название'] + [i.name for i in cinemas])

            # Записываем данные из словаря
            for row_num, (product, sales) in enumerate(data_dict.items(), start=1):
                worksheet.write(row_num, 0, product)  # Название продукта
                worksheet.write_row(row_num, 1, sales)  # Продажи

            # Создаем график
            chart = workbook.add_chart({'type': 'column'})

            # Добавляем данные к графику
            for i, product in enumerate(data_dict.keys()):
                chart.add_series({
                    'name': f'={worksheet.name}!$A${i + 1}',
                    'categories': f'={worksheet.name}!$B$1:$D$1',
                    'values': f'={worksheet.name}!$B${i + 1}:$D${i + 1}',
                })

            # Настраиваем график
            chart.set_title({'name': 'Загруженность кинотеатров'})
            chart.set_x_axis({'name': 'Кинотеатры'})
            chart.set_y_axis({'name': 'Количество сеансов'})

            # Вставляем график в лист
            worksheet.insert_chart('F2', chart)

            # Закрываем файл
            workbook.close()


main()

